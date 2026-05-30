#!/usr/bin/env python3
"""Reformat Stata code blocks into clean monospace 'code cards', fix mangled
comment/command lines, and add plain-English speaker notes."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "UMU 10-Non Linear Regression.pptx"
OUT = "UMU 10-Non Linear Regression.pptx"  # overwrite (backup made separately)

COMMENT = RGBColor(0x6A, 0x73, 0x7D)   # muted gray for * comments
CODE    = RGBColor(0x1F, 0x23, 0x28)   # near-black for commands
CARD_BG = RGBColor(0xF6, 0xF8, 0xFA)   # light code-card background
CARD_LN = RGBColor(0xD0, 0xD7, 0xDE)   # subtle border

def style_card(shape):
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_LN; shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    try: tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception: pass
    tf.margin_left = Pt(10); tf.margin_right = Pt(10)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)

def set_code(shape, code, size):
    style_card(shape)
    tf = shape.text_frame
    tf.clear()
    lines = code.split("\n")
    while lines and lines[0] == "": lines.pop(0)
    while lines and lines[-1] == "": lines.pop()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0; p.space_before = Pt(0); p.space_after = Pt(0)
        run = p.add_run(); run.text = ln
        f = run.font; f.name = "Consolas"; f.size = Pt(size)
        f.color.rgb = COMMENT if ln.lstrip().startswith("*") else CODE

# slide number (1-based) -> list of (shape_name, size, code)
SPEC = {
5: [("Content Placeholder 2", 18, """
* Load the data
sysuse auto, clear

* Explore the data
describe
list mpg weight in 1/5
"""), ("Content Placeholder 3", 18, """
* Variables
* - mpg    = miles per gallon   (outcome)
* - weight = vehicle weight, lb (predictor)

* Clinical parallel:
* drug concentration (mpg) falls as
* body mass (weight) rises
""")],
6: [("Content Placeholder 3", 17, """
sysuse auto, clear

twoway scatter mpg weight, ///
    title("MPG vs Vehicle Weight") ///
    xtitle("Weight (lbs)") ///
    ytitle("Miles per Gallon") ///
    scheme(s2color)
""")],
7: [("Rectangle 1", 18, """
* Step 2 - starting values, read off the plot
*   max MPG ~ 40   ->  a = 40
*   slow decay     ->  b = -0.0005

* Step 3 - fit the exponential decay model
nl (mpg = {a} * exp({b} * weight)), ///
    initial(a 40 b -0.0005)

* Same model, robust SEs (IRR-style interpretation)
nl (mpg = {a} * exp({b} * weight)), ///
    initial(a 40 b -0.0005) vce(robust)
""")],
8: [("Content Placeholder 2", 15, """
------------------------------------------------------------------
       mpg | Coefficient  Std. err.    t    P>|t|  [95% conf. int.]
-----------+------------------------------------------------------
        /a |   41.23456     2.34500  17.58  0.000   36.567  45.902
        /b |   -0.000478    0.00007  -6.83  0.000  -0.00062 -0.00034
------------------------------------------------------------------
""")],
10: [("Content Placeholder 2", 17, """
* Predict fitted values and residuals
predict fitted_exp, yhat
predict resid_exp, resid

* Diagnostic 1: residuals vs fitted
rvfplot, yline(0) ///
    title("Residuals vs Fitted - Exponential Model")
"""), ("Content Placeholder 3", 17, """
* Diagnostic 2: normality of residuals
qnorm resid_exp, ///
    title("Q-Q Plot of Residuals")

* Diagnostic 3: histogram
histogram resid_exp, normal ///
    title("Distribution of Residuals")
""")],
11: [("Rectangle 1", 18, """
* After fitting:
*   nl (mpg = {a}*exp({b}*weight)), initial(a 40 b -0.0005)

* Half-life = weight change that halves MPG
nlcom (half_life: ln(2) / -_b[/b])

* Output: half_life = 1450 lbs (95% CI: 1100, 1800)
* Meaning: every +1450 lb halves the MPG
""")],
12: [("Rectangle 1", 18, """
* Load the data
sysuse uslifeexp, clear

* Explore
list year le in 1/10

* Variables
* - year = calendar year       (predictor)
* - le   = life expectancy, yr (outcome)
""")],
14: [("Rectangle 1", 18, """
twoway scatter le year, ///
    title("US Life Expectancy Over Time") ///
    xtitle("Year") ytitle("Life Expectancy (years)")
""")],
16: [("Rectangle 2", 17, """
* Fit logistic growth model
nl (le = {asym} / (1 + exp({rate} * (year - {inflect})))), ///
    initial(asym 80 rate -0.1 inflect 1950)

* Store results for later comparison
estimates store logistic_model
""")],
17: [("Rectangle 4", 14, """
------------------------------------------------------------------
        le | Coefficient  Std. err.    t    P>|t|  [95% conf. int.]
-----------+------------------------------------------------------
     /asym |   79.234      1.234    64.21  0.000   76.745   81.723
     /rate |   -0.095      0.011    -8.64  0.000   -0.117   -0.073
  /inflect | 1958.45       2.101   931.78  0.000 1954.222 1962.678
------------------------------------------------------------------
""")],
18: [("Rectangle 1", 13, """
* Predict fitted values and residuals
predict fitted_log, yhat
predict resid_log, resid

* Actual vs predicted
twoway (scatter le year) ///
       (line fitted_log year, lpattern(dash) lcolor(red) lwidth(med)), ///
    title("Actual vs Predicted Life Expectancy") ///
    legend(label(1 "Actual") label(2 "Predicted"))

* Predict life expectancy for 2020
local asym = _b[/asym]
local rate = _b[/rate]
local inflect = _b[/inflect]
display "Predicted LE 2020: " `asym'/(1+exp(`rate'*(2020-`inflect')))

* Residuals over time
scatter resid_log year, yline(0) title("Residuals Over Time")
""")],
19: [("Rectangle 1", 15, """
* Simulate a dose-response study
clear
set seed 12345
set obs 40

* Dose levels: 0 to 100 mg/kg
gen dose = (_n - 1) * (100 / 39)

* True values: Vmax = 90%, Km (EC50) = 25 mg/kg
* response = (90*dose)/(25+dose) + noise
gen response = (90 * dose) / (25 + dose) + rnormal(0, 4)

* Plot the data
twoway scatter response dose, ///
    title("Dose-Response Curve") ///
    xtitle("Dose (mg/kg)") ytitle("Response (%)") ///
    ylabel(0(20)100)
""")],
21: [("Rectangle 1", 18, """
* Fit Michaelis-Menten
nl (response = {vmax} * dose / ({km} + dose)), ///
    initial(vmax 90 km 25)

* Same model with robust standard errors
nl (response = {vmax} * dose / ({km} + dose)), ///
    initial(vmax 90 km 25) vce(robust)
""")],
22: [("Content Placeholder 2", 15, """
------------------------------------------------------------------
  response | Coefficient  Std. err.    t    P>|t|  [95% conf. int.]
-----------+------------------------------------------------------
     /vmax |   89.456      2.345    38.15  0.000   84.712   94.200
       /km |   24.876      2.123    11.72  0.000   20.578   29.174
------------------------------------------------------------------
""")],
24: [("Rectangle 1", 17, """
* Predict fitted values and residuals
predict fitted_mm, yhat
predict resid_mm, resid

* Plot fitted curve over the data
twoway (scatter response dose) ///
       (line fitted_mm dose, lcolor(red) lwidth(med)), ///
    title("Michaelis-Menten Fit") ///
    xtitle("Dose (mg/kg)") ytitle("Response (%)")
""")],
25: [("Rectangle 1", 18, """
* Residual diagnostics
rvfplot, yline(0) title("Residuals vs Fitted")

* Cook's distance - find influential points
predict cooksd_mm, cooksd
scatter cooksd_mm _n, title("Cook's Distance") yline(4/40)
""")],
26: [("Rectangle 2", 15, """
* Simulate viral load after infection
clear
set seed 67890
set obs 20

* Days after infection
gen day = _n

* True values: start = 100, growth rate = 0.3
* viral_load = 100*exp(0.3*day) + noise
gen viral_load = 100 * exp(0.3 * day) + rnormal(0, 200)

* Plot the data
twoway scatter viral_load day, ///
    title("Viral Load Over Time") ///
    xtitle("Days Post-Infection") ///
    ytitle("Viral Load (copies/mL)")
""")],
27: [("Rectangle 1", 14, """
* Starting values: a = 100, b = 0.3
nl (viral_load = {a} * exp({b} * day)), ///
    initial(a 100 b 0.3)

------------------------------------------------------------------
viral_load | Coefficient  Std. err.    t    P>|t|  [95% conf. int.]
-----------+------------------------------------------------------
        /a |   98.456     12.34     7.98  0.000   73.456  123.456
        /b |   0.312       0.045    6.93  0.000    0.218    0.406
------------------------------------------------------------------

* Calculate doubling time
nlcom (doubling_time: ln(2) / _b[/b])
* Doubling time = 2.22 days (95% CI: 1.71, 2.73)
""")],
28: [("Rectangle 1", 16, """
* Linear model on raw data
regress mpg weight
estimates store linear_model

* Non-linear model (exponential decay)
nl (mpg = {a} * exp({b} * weight)), initial(a 40 b -0.0005)
estimates store nl_model

* Compare AIC / BIC
estimates stats linear_model nl_model

* Compare residual spread
predict resid_linear, resid
summarize resid_linear resid_exp
""")],
30: [("Rectangle 1", 18, """
* Run this after EVERY non-linear regression

* 1. Check convergence (should be 1)
di e(converged)

* 2. Residuals vs fitted
rvfplot, yline(0)

* 3. Normality of residuals
qnorm resid
""")],
31: [("Rectangle 1", 16, """
* 3. Normality of residuals
qnorm resid
swilk resid        // Shapiro-Wilk: p > 0.05 = normal

* 4. Influential points
predict cooksd, cooksd
summarize cooksd, detail

* 5. Parameter plausibility
*    do the estimates make clinical sense?

* 6. Profile-likelihood CIs (better for small n)
nl (y = {a} * exp({b} * x)), initial(a 10 b -0.1) level(95)
""")],
33: [("Rectangle 1", 17, """
* Trace each iteration
nl (y = {a} * exp({b} * x)), initial(a 40 b -0.0005) trace

* Increase the iteration limit
nl (y = {a} * exp({b} * x)), initial(a 40 b -0.0005) iterate(100)

* Try a different algorithm
nl (y = {a} * exp({b} * x)), initial(a 40 b -0.0005) technique(gn)
""")],
36: [("Rectangle 1", 16, """
* ================================================
* EXERCISE: Hospital infection-rate data
* ================================================
clear
set seed 999
set obs 30

* Hospital occupancy (%)
gen occupancy = 40 + (_n - 1) * (95 - 40) / 29

* True logistic parameters:
*   asymptote (max rate) = 12 per 1000 patient-days
*   inflection point     = 75% occupancy
""")],
37: [("Rectangle 1", 14, """
* Growth rate = 0.3
gen infection_rate = 12 / (1 + exp(0.3*(occupancy-75))) + rnormal(0, 0.5)

* Round for realism
replace infection_rate = round(infection_rate, 0.1)

* Label variables
label variable occupancy "Hospital Bed Occupancy (%)"
label variable infection_rate "Infection Rate (per 1000 patient-days)"

* View and plot
list in 1/10
twoway scatter infection_rate occupancy, ///
    title("Hospital Infection Rate by Occupancy") ///
    xtitle("Occupancy (%)") ///
    ytitle("Infection Rate (per 1000 patient-days)")
""")],
38: [("Rectangle 1", 16, """
nl (infection_rate = {asym} / ///
    (1 + exp({rate_g}*(occupancy-{inflect})))), ///
    initial(asym 12 rate_g 0.3 inflect 75)
""")],
39: [("Rectangle 1", 17, """
* Predict fitted values and residuals
predict fitted_rate, yhat
predict resid_ex, resid

* Plot actual vs predicted
twoway (scatter infection_rate occupancy) ///
       (line fitted_rate occupancy, lcolor(red)), ///
    title("Model Fit")
""")],
40: [("Rectangle 1", 16, """
* For a logistic model the inflection point IS the EC50.
* Retrieve it after fitting:
nl (infection_rate = {asym} / ///
    (1 + exp({rate_g}*(occupancy-{inflect})))), ///
    initial(asym 12 rate_g 0.3 inflect 75)
""")],
41: [("Rectangle 1", 16, """
* Display the inflection point with a 95% CI
di "Occupancy at 50% max rate: " _b[/inflect] " (95% CI: " ///
   _b[/inflect] - 1.96*_se[/inflect] ", " ///
   _b[/inflect] + 1.96*_se[/inflect] ")"

* Alternative using nlcom
nlcom (ec50: _b[/inflect])

* Interpretation: above this occupancy, the
* infection rate accelerates rapidly
""")],
42: [("Rectangle 1", 15, """
* ================================================
* EXERCISE SOLUTION
* ================================================
* Tasks 1-2: create and plot data (see earlier slides)

* Tasks 3-4: fit the logistic model
nl (infection_rate = {asym} / ///
    (1 + exp({rate_g}*(occupancy-{inflect})))), ///
    initial(asym 12 rate_g 0.3 inflect 75)

* Task 5: interpret parameters
*   /asym    = maximum infection rate (ceiling)
*   /rate_g  = steepness of the rise
*   /inflect = occupancy at 50% of maximum
""")],
43: [("Rectangle 1", 14, """
* Task 6: predict and plot
predict fitted_rate, yhat
twoway (scatter infection_rate occupancy) ///
       (line fitted_rate occupancy, lcolor(red) lwidth(med)), ///
    title("Actual vs Predicted Infection Rate") ///
    xtitle("Occupancy (%)") ytitle("Infection Rate (per 1000)")

* Task 7: EC50 (= /inflect)
nlcom (occupancy_50pct: _b[/inflect])

* Task 8: residual diagnostics
predict resid_ex, resid
qnorm resid_ex, title("Q-Q Plot of Residuals")
rvfplot, yline(0) title("Residuals vs Fitted")
""")],
44: [("Rectangle 1", 13, """
------------------------------------------------------------------
infection_~e | Coefficient  Std. err.   t    P>|t|  [95% conf. int.]
-------------+----------------------------------------------------
       /asym |   12.234     0.345   35.46  0.000   11.527  12.941
     /rate_g |    0.312     0.045    6.93  0.000    0.220   0.404
    /inflect |   75.123     0.876   85.76  0.000   73.345  76.901
------------------------------------------------------------------

Occupancy at 50% max: 75.1% (95% CI: 73.3%, 76.9%)
""")],
}

NOTES = {
5: "We borrow Stata's built-in 'auto' dataset as a stand-in for a pharmacokinetic curve: MPG plays the role of drug concentration, vehicle weight the role of body mass. sysuse loads it; describe and list let us eyeball the variables before modelling.",
6: "Always plot before fitting. The scatter of MPG against weight shows a curved, decaying relationship - that shape is what tells us an exponential decay model (not a straight line) is appropriate. The /// just continues one long command across several lines.",
7: "nl fits a non-linear model. The {a} and {b} in curly braces are the parameters Stata will estimate. initial() gives it sensible starting guesses read off the plot (max MPG ~40, gentle decay). The second version adds vce(robust) for standard errors that tolerate non-constant variance.",
8: "Read the two parameters: /a is the ceiling (MPG at zero weight, ~41), /b is the decay rate (negative, so MPG falls as weight rises). Both p-values are 0.000, so both are clearly different from zero.",
10: "After any fit, predict pulls out the fitted values (yhat) and residuals. The three plots check the fit: residuals vs fitted should show no pattern, the Q-Q plot and histogram check that residuals are roughly normal.",
11: "nlcom computes a NEW quantity from the fitted parameters, with a correct confidence interval. Here half-life = ln(2)/-b answers 'how much extra weight halves the MPG?' - the same idea as a drug's half-life.",
12: "Switching datasets: US life expectancy 1900-2000, which rises in an S-shape - our analogy for tumour growth or CD4 recovery. Just load and inspect here. (The original 'keep if region' line was dropped - this dataset has no region variable, so it would error.)",
14: "Plot first again. The S-shaped rise of life expectancy over time is the signature of logistic growth: slow, then fast, then levelling off at a ceiling.",
16: "The logistic model has three meaningful parameters: asym (the ceiling), rate (how steep), inflect (the year of fastest growth). estimates store saves the fit so we can compare models later.",
17: "asym ~79 years is the estimated ceiling, inflect ~1958 is the year growth was fastest, rate sets the steepness. Tight confidence intervals and tiny p-values mean all three are precisely estimated.",
18: "predict gets fitted values; the twoway overlays the fitted curve on the raw data. We then store the three parameters in locals and plug them into the logistic formula to forecast 2020. Finally we plot residuals over time to check for leftover pattern.",
19: "Here we simulate clean dose-response data so we know the truth. set seed makes it reproducible. We build doses from 0-100 mg/kg, generate responses from a known Michaelis-Menten curve (Vmax=90, EC50=25) plus noise, then plot.",
21: "Michaelis-Menten: response rises fast at low dose then saturates. vmax is the maximum response, km is the dose giving half of it (the EC50). The second run adds robust standard errors.",
22: "vmax ~89% is the saturating maximum; km ~24.9 mg/kg is the EC50 - the dose that achieves half the maximum effect. Both are estimated precisely.",
24: "Standard post-fit check: predict the fitted curve and residuals, then overlay the smooth fitted line on the scatter to see how well the curve tracks the data.",
25: "Two more diagnostics: rvfplot looks for pattern in the residuals (there should be none), and Cook's distance flags any single point that unduly influences the fit. The 4/40 line is the common cut-off (4/n).",
26: "Exponential GROWTH this time - viral load climbing after infection. Same recipe: seed, simulate from a known curve (start 100, growth 0.3) with noise, then plot the rise.",
27: "Same nl machinery, now fitting growth. /a recovers the starting load (~98), /b the growth rate (~0.31). nlcom then turns the rate into a doubling time = ln(2)/b ~ 2.2 days - the clinically intuitive number.",
28: "Why bother with non-linear? Fit both a straight line and the exponential model to the same data, store each, then compare AIC/BIC (lower = better) and residual spread. Here non-linear wins modestly - so choose on clinical interpretability.",
30: "A reusable checklist to run after every non-linear fit. First: e(converged) must equal 1, or nothing else is trustworthy. Then the residual and normality plots.",
31: "Continuing the checklist: swilk is a formal normality test (p>0.05 = normal), Cook's distance finds influential points, and level(95) requests profile-likelihood confidence intervals, which are more honest than the default ones in small samples.",
33: "When a model won't converge: trace shows each iteration so you can see where it goes wrong; iterate(100) gives it more attempts; technique(gn) switches to the Gauss-Newton algorithm. Often the real fix is better initial() values.",
36: "The hands-on exercise. Build a hospital dataset where infection rate follows a logistic S-curve against bed occupancy: low until ~75% full, then rising sharply. This slide just creates occupancy and states the true parameters.",
37: "Generate the infection rate from the logistic formula plus noise, round it, label the variables, then list and plot. (Straight quotes are used here so the commands paste and run cleanly.)",
38: "Hint: the model structure to fit. Note the parameters in braces and the starting values - asymptote ~12 (the plot's ceiling), positive growth rate, inflection ~75%.",
39: "After fitting, predict the fitted curve and residuals and overlay the fit on the data to judge how well the S-curve matches.",
40: "A neat property: for a logistic curve the inflection point IS the EC50 - the occupancy at 50% of the maximum infection rate. So you read it straight off the /inflect parameter.",
41: "Two ways to report the EC50 with a confidence interval: di builds the sentence by hand using the stored coefficient and its standard error; nlcom does it more cleanly and automatically.",
42: "Worked solution, part 1: create/plot the data, fit the logistic model, and interpret each parameter - ceiling, steepness, and the 50% point.",
43: "Worked solution, part 2: predict and plot the fit, report the EC50 via nlcom, then run the residual diagnostics (Q-Q and residual-vs-fitted) to confirm the model is sound.",
44: "Expected results. The ceiling is ~12.2 infections per 1000 patient-days, and the 50% point sits at ~75% occupancy - the actionable threshold: keep occupancy below 75% and infection rates stay low.",
}

prs = Presentation(SRC)
applied, noted = 0, 0
for snum, items in SPEC.items():
    slide = prs.slides[snum - 1]
    byname = {sh.name: sh for sh in slide.shapes}
    for shname, size, code in items:
        sh = byname.get(shname)
        if sh is None or not sh.has_text_frame:
            print(f"!! slide {snum}: shape {shname!r} not found"); continue
        set_code(sh, code, size)
        applied += 1
for snum, note in NOTES.items():
    slide = prs.slides[snum - 1]
    slide.notes_slide.notes_text_frame.text = note
    noted += 1

prs.save(OUT)
print(f"Applied {applied} code cards, {noted} note blocks. Saved {OUT}")
